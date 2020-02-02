#!/usr/bin/env python

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
import os

IMG_TYPES = ['png', 'jpg', 'jpeg']

def listdir(path, list_name):
  for file in os.listdir(path):  
    file_path = os.path.join(path, file)  
    if os.path.isdir(file_path):  
      list_name = listdir(file_path, list_name)  
    else:  
      list_name.append(file_path)
  return list_name

def add_text_to_slide(slide, text):
  txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
  p = txBox.text_frame.add_paragraph()
  p.text = text
  p.font.color.rgb = RGBColor(255, 0, 0)

if __name__ == "__main__":
  # 获取所有图片
  picture_list = listdir('./pictures/', [])
  img_list = []
  for item in picture_list:
    if item.split('.')[-1].lower() in IMG_TYPES:
      img_list.append(item)

  # 获取所有音频
  sound_list = listdir('./mp3/', [])
  print(sound_list)

  # 获取所有视频
  video_list = listdir('./mp4/', [])
  print(video_list)

  prs = Presentation()
  slide_list = []

  for img in img_list:
    # prs.slide_layouts中一共预存有1-48种，采用第六种为空白幻灯片
    # 创建一张空白幻灯片
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # 左上角0，0位置
    left = top = Inches(0)
    # 设置长宽自适应ppt大小
    height = prs.slide_height
    width = prs.slide_width

    # 添加一张图片
    slide.shapes.add_picture(img, left, top, width=width, height=height)
    slide_list.append(slide)

  # 添加一个没有封面的音频
  slide_list[0].shapes.add_movie(sound_list[0], Inches(0), Inches(0), Inches(0.5), Inches(0.25))
  add_text_to_slide(slide_list[0], '添加了一张图片、一个没有封面的音频')

  # 添加一个有默认封面的音频
  slide_list[1].shapes.add_movie(sound_list[1], Inches(2), Inches(0), Inches(0.5), Inches(0.5), poster_frame_image='./sound_icon.png')
  add_text_to_slide(slide_list[1], '添加了一张图片、一个有默认封面的音频')

  # 添加一个没有默认封面的视频
  slide_list[2].shapes.add_movie(video_list[1], Inches(0), Inches(3), Inches(3), Inches(2))
  add_text_to_slide(slide_list[2], '添加了一张图片、一个没有封面的视频')

  # 添加一个有默认封面的视频
  slide_list[3].shapes.add_movie(video_list[0], Inches(3), Inches(3), Inches(3), Inches(2), poster_frame_image='./video_icon.png')
  add_text_to_slide(slide_list[3], '添加了一张图片、一个有默认封面的视频')


  prs.save('test.pptx')